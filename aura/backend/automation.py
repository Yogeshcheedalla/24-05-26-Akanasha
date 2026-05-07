import asyncio
import os
import re
import shutil
import subprocess
import tempfile
import time
from pathlib import Path
from urllib.parse import quote_plus

try:
    import winreg
except ImportError:  # pragma: no cover - Windows-only helper
    winreg = None

import fitz
import pyautogui
import pygetwindow as gw
from pywinauto import Desktop
from pptx import Presentation
from pptx.util import Inches

pyautogui.FAILSAFE = False


APP_LAUNCH_COMMANDS: dict[str, list[str]] = {
    "notepad": ["notepad.exe"],
    "calculator": ["calc.exe"],
    "calc": ["calc.exe"],
    "file explorer": ["explorer.exe"],
    "explorer": ["explorer.exe"],
    "vscode": ["Code.exe"],
    "visual studio code": ["Code.exe"],
    "chrome": [
        "chrome.exe",
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
    ],
    "google chrome": [
        "chrome.exe",
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
    ],
    "brave": [
        "brave.exe",
        r"C:\Program Files\BraveSoftware\Brave-Browser\Application\brave.exe",
        r"C:\Program Files (x86)\BraveSoftware\Brave-Browser\Application\brave.exe",
    ],
    "brave browser": [
        "brave.exe",
        r"C:\Program Files\BraveSoftware\Brave-Browser\Application\brave.exe",
        r"C:\Program Files (x86)\BraveSoftware\Brave-Browser\Application\brave.exe",
    ],
    "edge": [
        "start msedge",
        "microsoft-edge:",
        "msedge.exe",
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ],
    "microsoftedge": [
        "start msedge",
        "microsoft-edge:",
        "msedge.exe",
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ],
    "microsoft edge": [
        "start msedge",
        "microsoft-edge:",
        "msedge.exe",
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ],
    "command prompt": ["cmd.exe"],
    "cmd": ["cmd.exe"],
    "powershell": ["powershell.exe"],
    "whatsapp": [
        "shell:AppsFolder\\5319275A.WhatsAppDesktop_cv1g1gvanyjgm!App",
        "whatsapp:",
        "WhatsApp.exe",
    ],
    "whatsapp desktop": [
        "shell:AppsFolder\\5319275A.WhatsAppDesktop_cv1g1gvanyjgm!App",
        "whatsapp:",
        "WhatsApp.exe",
    ],
    "telegram": ["tg:", "Telegram.exe"],
    "telegram desktop": ["tg:", "Telegram.exe"],
    "discord": ["discord:", "Discord.exe"],
    "discord desktop": ["discord:", "Discord.exe"],
    "word": ["winword.exe"],
    "excel": ["excel.exe"],
    "powerpoint": ["powerpnt.exe"],
    "settings": ["ms-settings:"],
    "windows settings": ["ms-settings:"],
    "terminal": ["wt.exe"],
    "windows terminal": ["wt.exe"],
    "control panel": ["control.exe"],
}

START_MENU_FRIENDLY_NAMES: dict[str, str] = {
    "notepad": "Notepad",
    "calculator": "Calculator",
    "calc": "Calculator",
    "file explorer": "File Explorer",
    "explorer": "File Explorer",
    "vscode": "Visual Studio Code",
    "visual studio code": "Visual Studio Code",
    "chrome": "Google Chrome",
    "google chrome": "Google Chrome",
    "brave": "Brave",
    "brave browser": "Brave",
    "edge": "Microsoft Edge",
    "microsoftedge": "Microsoft Edge",
    "microsoft edge": "Microsoft Edge",
    "command prompt": "Command Prompt",
    "cmd": "Command Prompt",
    "powershell": "PowerShell",
    "whatsapp": "WhatsApp",
    "whatsapp desktop": "WhatsApp",
    "telegram": "Telegram",
    "telegram desktop": "Telegram",
    "discord": "Discord",
    "discord desktop": "Discord",
    "word": "Word",
    "excel": "Excel",
    "powerpoint": "PowerPoint",
    "settings": "Settings",
    "windows settings": "Settings",
    "terminal": "Windows Terminal",
    "windows terminal": "Windows Terminal",
    "control panel": "Control Panel",
}

APP_REGISTRY_EXECUTABLES: dict[str, list[str]] = {
    "edge": ["msedge.exe"],
    "microsoftedge": ["msedge.exe"],
    "microsoft edge": ["msedge.exe"],
    "chrome": ["chrome.exe"],
    "google chrome": ["chrome.exe"],
    "brave": ["brave.exe"],
    "brave browser": ["brave.exe"],
    "word": ["winword.exe"],
    "excel": ["excel.exe"],
    "powerpoint": ["powerpnt.exe"],
    "whatsapp": ["WhatsApp.exe"],
    "whatsapp desktop": ["WhatsApp.exe"],
    "telegram": ["Telegram.exe"],
    "telegram desktop": ["Telegram.exe"],
    "discord": ["Discord.exe"],
    "discord desktop": ["Discord.exe"],
}

APP_NAME_ALIASES: dict[str, str] = {
    "ms edge": "microsoft edge",
    "microsoftedge": "microsoft edge",
    "edge browser": "microsoft edge",
    "browser edge": "microsoft edge",
    "google chrome": "chrome",
    "chrome browser": "chrome",
    "brave browser": "brave",
    "whats app": "whatsapp",
    "whatsup": "whatsapp",
    "watsup": "whatsapp",
    "whatsap": "whatsapp",
    "whatsapp desktop": "whatsapp",
    "telegram desktop": "telegram",
    "discord desktop": "discord",
    "windows settings": "settings",
    "windows terminal": "terminal",
    "whatsup desktop": "whatsapp",
    "whatsapp web": "whatsapp",
    "microsoft edge desktop": "microsoft edge",
    "edge desktop": "microsoft edge",
}

AMMA_ONLY_CONTACT_ALIASES: dict[str, str] = {
    "amma": "Amma",
    "ammaa": "Amma",
    "ammah": "Amma",
    "amm": "Amma",
    "am ma": "Amma",
    "mummy": "Amma",
    "mumyy": "Amma",
    "mommy": "Amma",
    "mom": "Amma",
    "mother": "Amma",
    "mamma": "Amma",
    "mumy": "Amma",
}

WHATSAPP_CONTACT_DISPLAY_NAMES: dict[str, list[str]] = {
    "Amma": ["Amma💗", "Amma"],
}

SPECIAL_FOLDERS: dict[str, Path] = {
    "downloads": Path.home() / "Downloads",
    "desktop": Path.home() / "Desktop",
    "documents": Path.home() / "Documents",
    "pictures": Path.home() / "Pictures",
    "videos": Path.home() / "Videos",
    "music": Path.home() / "Music",
}


def _lookup_windows_app_path(executable_name: str) -> str | None:
    if not winreg or not executable_name.lower().endswith(".exe"):
        return None

    registry_path = fr"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\{executable_name}"
    for hive in (winreg.HKEY_CURRENT_USER, winreg.HKEY_LOCAL_MACHINE):
        try:
            with winreg.OpenKey(hive, registry_path) as key:
                path, _ = winreg.QueryValueEx(key, None)
                if path and os.path.exists(path):
                    return path
        except OSError:
            continue
    return None


def _normalize_app_name(app_name: str) -> str:
    normalized = re.sub(r"\s+", " ", app_name.strip().lower())
    normalized = normalized.replace("whats up", "whatsapp")
    normalized = normalized.replace("whatsup", "whatsapp")
    normalized = normalized.replace("watsup", "whatsapp")
    normalized = normalized.replace("whastapp", "whatsapp")
    normalized = APP_NAME_ALIASES.get(normalized, normalized)
    return normalized


def _normalize_amma_only_contact(contact: str) -> str:
    normalized = re.sub(r"\s+", " ", contact.strip().lower())
    normalized = re.sub(r"\b(?:on|in|through|via)\s+whatsapp(?:\s+desktop|\s+app|\s+website|\s+web)?\b", "", normalized)
    normalized = normalized.replace("whatsapp", "").replace("desktop", "").replace("website", "").replace("web", "").strip()
    normalized = re.sub(r"[^a-z\s]", " ", normalized)
    normalized = re.sub(r"\s+", " ", normalized).strip()
    normalized = AMMA_ONLY_CONTACT_ALIASES.get(normalized, normalized.title())
    if normalized != "Amma":
        raise ValueError(
            "For safety, this build only allows WhatsApp message automation to the Amma contact."
        )
    return normalized


def _get_whatsapp_display_target(contact: str) -> str:
    if contact == "Amma":
        return "Amma"
    display_options = WHATSAPP_CONTACT_DISPLAY_NAMES.get(contact, [contact])
    return display_options[0]


def _get_whatsapp_search_query(contact: str) -> str:
    if contact == "Amma":
        return "Amma"
    return re.sub(r"[^\w\s]", "", contact).strip() or contact


def _normalize_visible_text(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", (value or "").strip().lower())


def _get_window_bounds(window) -> tuple[int, int, int, int]:
    if hasattr(window, "rectangle"):
        rect = window.rectangle()
        return (
            int(rect.left),
            int(rect.top),
            int(rect.right - rect.left),
            int(rect.bottom - rect.top),
        )
    return int(window.left), int(window.top), int(window.width), int(window.height)


def _wait_for_whatsapp_window(timeout_seconds: float = 12.0):
    deadline = time.time() + timeout_seconds
    last_error: Exception | None = None
    while time.time() < deadline:
        try:
            windows = []
            for window in Desktop(backend="uia").windows():
                try:
                    title = (window.window_text() or "").strip()
                    class_name = (window.element_info.class_name or "").strip()
                except Exception:
                    continue
                if class_name != "WinUIDesktopWin32WindowClass":
                    continue
                if not title or not title.lower().endswith("whatsapp"):
                    continue
                windows.append(window)
            if windows:
                windows.sort(
                    key=lambda candidate: (
                        0 if (candidate.window_text() or "").strip().lower() == "whatsapp" else 1,
                        -((_get_window_bounds(candidate)[2]) * (_get_window_bounds(candidate)[3])),
                    )
                )
                window = windows[0]
                try:
                    if hasattr(window, "restore"):
                        window.restore()
                except Exception:
                    pass
                try:
                    window.set_focus()
                except Exception:
                    try:
                        window.click_input()
                    except Exception:
                        pass
                return window
        except Exception as exc:
            last_error = exc
        time.sleep(0.35)
    raise RuntimeError(f"WhatsApp window did not become ready in time: {last_error}")


def _click_window_ratio(window, x_ratio: float, y_ratio: float) -> tuple[int, int]:
    left, top, width, height = _get_window_bounds(window)
    if width <= 0 or height <= 0:
        raise RuntimeError("WhatsApp window does not have a usable size yet.")
    x = left + int(width * x_ratio)
    y = top + int(height * y_ratio)
    pyautogui.click(x, y)
    return x, y


def _click_element_center(element) -> bool:
    try:
        rect = element.rectangle()
        x = int((rect.left + rect.right) / 2)
        y = int((rect.top + rect.bottom) / 2)
        pyautogui.click(x, y)
        return True
    except Exception:
        return False


def _get_element_rectangle(element):
    try:
        return element.rectangle()
    except Exception:
        return None


def _clear_active_field() -> None:
    pyautogui.hotkey("ctrl", "a")
    pyautogui.sleep(0.08)
    pyautogui.press("backspace")


def _iter_descendants(window):
    try:
        return list(window.descendants())
    except Exception:
        return []


def _iter_element_ancestors(element, max_depth: int = 6):
    current = element
    seen_ids: set[int] = set()
    for _ in range(max_depth):
        if current is None:
            break
        object_id = id(current)
        if object_id in seen_ids:
            break
        seen_ids.add(object_id)
        yield current
        try:
            current = current.parent()
        except Exception:
            break


def _is_left_results_pane(window, rect) -> bool:
    if rect is None:
        return False
    left, top, width, height = _get_window_bounds(window)
    if width <= 0 or height <= 0:
        return False
    x_center = (rect.left + rect.right) / 2
    y_center = (rect.top + rect.bottom) / 2
    x_min = left + (width * 0.03)
    x_max = left + (width * 0.38)
    y_min = top + (height * 0.18)
    y_max = top + (height * 0.86)
    return x_min <= x_center <= x_max and y_min <= y_center <= y_max


def _find_whatsapp_contact_element(window, display_options: list[str]):
    normalized_targets = [_normalize_visible_text(option) for option in display_options if option]
    best_match = None
    best_score = -1

    for element in _iter_descendants(window):
        try:
            text = (element.window_text() or "").strip()
            if not text:
                continue
            normalized_text = _normalize_visible_text(text)
            if not normalized_text:
                continue
            control_type = (element.element_info.control_type or "").lower()
            rect = _get_element_rectangle(element)
        except Exception:
            continue

        if not _is_left_results_pane(window, rect):
            continue

        for target in normalized_targets:
            if not target:
                continue
            if normalized_text == target:
                score = 140
            elif target in normalized_text or normalized_text in target:
                score = 80
            else:
                continue

            if control_type in {"listitem", "dataitem"}:
                score += 35
            elif control_type == "text":
                score += 18
            if rect is not None:
                row_width = rect.right - rect.left
                row_height = rect.bottom - rect.top
                if row_width > 120:
                    score += 12
                if row_height > 24:
                    score += 8
            if best_match is None or score > best_score:
                best_match = element
                best_score = score

    return best_match


def _pick_whatsapp_click_target(window, element):
    best_candidate = None
    best_score = -1

    for candidate in _iter_element_ancestors(element):
        rect = _get_element_rectangle(candidate)
        if not _is_left_results_pane(window, rect):
            continue
        try:
            control_type = (candidate.element_info.control_type or "").lower()
        except Exception:
            control_type = ""
        row_width = (rect.right - rect.left) if rect is not None else 0
        row_height = (rect.bottom - rect.top) if rect is not None else 0
        score = 0
        if control_type in {"listitem", "dataitem"}:
            score += 50
        elif control_type == "pane":
            score += 25
        elif control_type == "text":
            score += 10
        score += min(int(row_width / 20), 30)
        score += min(int(row_height / 10), 20)
        if score > best_score:
            best_candidate = candidate
            best_score = score

    return best_candidate or element


async def _wait_for_whatsapp_contact_match(window, display_options: list[str], timeout_seconds: float = 3.2):
    deadline = time.time() + timeout_seconds
    while time.time() < deadline:
        match = _find_whatsapp_contact_element(window, display_options)
        if match is not None:
            return match
        await asyncio.sleep(0.12)
    return None


def _is_whatsapp_chat_open(window, display_options: list[str]) -> bool:
    normalized_targets = {_normalize_visible_text(option) for option in display_options if option}
    left, top, width, height = _get_window_bounds(window)
    if width <= 0 or height <= 0:
        return False
    x_min = left + (width * 0.22)
    x_max = left + (width * 0.70)
    y_min = top + (height * 0.02)
    y_max = top + (height * 0.18)

    for element in _iter_descendants(window):
        try:
            text = (element.window_text() or "").strip()
            if not text:
                continue
            normalized_text = _normalize_visible_text(text)
            if not normalized_text:
                continue
            rect = _get_element_rectangle(element)
            if rect is None:
                continue
        except Exception:
            continue

        x_center = (rect.left + rect.right) / 2
        y_center = (rect.top + rect.bottom) / 2
        if not (x_min <= x_center <= x_max and y_min <= y_center <= y_max):
            continue

        for target in normalized_targets:
            if normalized_text == target or target in normalized_text:
                return True
    return False


async def _open_whatsapp_contact_guarded(contact: str) -> dict[str, str]:
    cleaned_contact = _normalize_amma_only_contact(contact)
    display_options = WHATSAPP_CONTACT_DISPLAY_NAMES.get(cleaned_contact, [cleaned_contact])
    display_target = _get_whatsapp_display_target(cleaned_contact)
    search_query = _get_whatsapp_search_query(cleaned_contact)

    _launch_windows_app("whatsapp")
    await asyncio.sleep(1.2)
    window = _wait_for_whatsapp_window()
    await asyncio.sleep(0.18)

    # Focus the left search box first; fall back to Ctrl+F if needed.
    _click_window_ratio(window, 0.13, 0.14)
    await asyncio.sleep(0.14)
    pyautogui.hotkey("ctrl", "f")
    await asyncio.sleep(0.14)
    _clear_active_field()
    await asyncio.sleep(0.05)
    pyautogui.write(search_query, interval=0.02)

    match = await _wait_for_whatsapp_contact_match(window, display_options)
    if match is None:
        raise RuntimeError(
            f"WhatsApp did not show a safe contact match for '{display_target}'. Refusing to send."
        )

    clicked = False
    click_target = _pick_whatsapp_click_target(window, match)
    for candidate in (click_target, match):
        if candidate is None:
            continue
        try:
            candidate.click_input()
            clicked = True
            break
        except Exception:
            if _click_element_center(candidate):
                clicked = True
                break

    if not clicked:
        raise RuntimeError(
            f"WhatsApp found '{display_target}' but could not safely click the result row."
        )

    await asyncio.sleep(0.35)
    if not _is_whatsapp_chat_open(window, display_options):
        pyautogui.press("enter")
        await asyncio.sleep(0.35)

    if not _is_whatsapp_chat_open(window, display_options):
        raise RuntimeError(
            f"WhatsApp did not open the verified '{display_target}' chat header after selecting the result."
        )

    return {
        "contact": cleaned_contact,
        "display_target": display_target,
        "search_query": search_query,
    }


def _resolve_launch_commands(app_name: str, commands: list[str]) -> list[str]:
    resolved: list[str] = []
    seen: set[str] = set()

    for command in commands:
        candidate = command
        lower_candidate = command.lower()
        if not (lower_candidate.startswith("shell:") or re.match(r"^[a-z][a-z0-9+.-]*:$", lower_candidate)):
            registry_match = _lookup_windows_app_path(Path(command).name)
            if registry_match:
                candidate = registry_match
        key = candidate.lower()
        if key not in seen:
            seen.add(key)
            resolved.append(candidate)

    for executable_name in APP_REGISTRY_EXECUTABLES.get(app_name, []):
        registry_match = _lookup_windows_app_path(executable_name)
        if registry_match and registry_match.lower() not in seen:
            seen.add(registry_match.lower())
            resolved.append(registry_match)

    return resolved


def _launch_windows_app(app_name: str) -> str:
    normalized = _normalize_app_name(app_name)
    if not normalized:
        raise ValueError("An app name is required.")

    commands = _resolve_launch_commands(normalized, APP_LAUNCH_COMMANDS.get(normalized, [normalized]))
    search_label = START_MENU_FRIENDLY_NAMES.get(normalized, app_name.strip())
    last_error: Exception | None = None

    for command in commands:
        try:
            lower_command = command.lower()
            if lower_command.startswith("start "):
                start_target = command[6:].strip()
                subprocess.Popen(
                    ["cmd.exe", "/c", "start", "", start_target],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                )
                return normalized
            if lower_command.startswith("shell:") or re.match(r"^[a-z][a-z0-9+.-]*:$", lower_command):
                if lower_command.startswith("shell:"):
                    subprocess.Popen(
                        ["explorer.exe", command],
                        stdout=subprocess.DEVNULL,
                        stderr=subprocess.DEVNULL,
                    )
                else:
                    if hasattr(os, "startfile"):
                        os.startfile(command)  # type: ignore[attr-defined]
                    else:
                        subprocess.Popen(
                            ["cmd.exe", "/c", "start", "", command],
                            stdout=subprocess.DEVNULL,
                            stderr=subprocess.DEVNULL,
                        )
            else:
                resolved = shutil.which(command) or command
                subprocess.Popen([resolved], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            return normalized
        except Exception as exc:  # pragma: no cover - best effort on host desktop
            last_error = exc

    # Fall back to the Windows Start menu search so friendly names like
    # "Word", "Settings", or installed local apps can still launch.
    try:
        pyautogui.press("win")
        pyautogui.sleep(0.5)
        pyautogui.write(search_label, interval=0.03)
        pyautogui.sleep(0.6)
        pyautogui.press("enter")
        return normalized
    except Exception as exc:  # pragma: no cover - best effort on host desktop
        last_error = exc

    raise RuntimeError(f"Could not open app '{app_name}': {last_error}")


async def _whatsapp_open_contact_and_send(contact: str, message: str) -> dict[str, str | bool]:
    cleaned_message = message.strip()
    if not cleaned_message:
        raise ValueError("A message is required for WhatsApp automation.")

    contact_payload = await _open_whatsapp_contact_guarded(contact)
    window = _wait_for_whatsapp_window()

    # Click the composer separately so the message never bleeds into the search field.
    _click_window_ratio(window, 0.57, 0.955)
    await asyncio.sleep(0.28)
    pyautogui.write(cleaned_message, interval=0.025)
    await asyncio.sleep(0.12)
    pyautogui.press("enter")

    return {
        "success": True,
        "message": (
            f"Opened WhatsApp, selected {contact_payload['contact']} using the guarded search flow, "
            f"and sent the message."
        ),
        "note": (
            f"Safety guard used the exact WhatsApp search target '{contact_payload['display_target']}' "
            "before clicking the result row and the message composer."
        ),
    }


async def _whatsapp_open_contact(contact: str) -> dict[str, str | bool]:
    contact_payload = await _open_whatsapp_contact_guarded(contact)

    return {
        "success": True,
        "message": f"Opened WhatsApp and focused the chat for {contact_payload['contact']}.",
        "note": f"Used the exact search target '{contact_payload['display_target']}' before opening the chat.",
    }


def _launch_windows_app_with_url(app_name: str, target_url: str) -> tuple[str, str]:
    normalized = _normalize_app_name(app_name)
    if not normalized:
        raise ValueError("A browser app name is required.")

    url = target_url.strip()
    if not url:
        raise ValueError("A website URL is required.")
    if not re.match(r"^https?://", url, flags=re.IGNORECASE):
        url = f"https://{url}"

    commands = _resolve_launch_commands(normalized, APP_LAUNCH_COMMANDS.get(normalized, [normalized]))
    last_error: Exception | None = None

    for command in commands:
        lower_command = command.lower()
        if lower_command.startswith("shell:"):
            continue
        if lower_command.startswith("start "):
            start_target = command[6:].strip()
            try:
                subprocess.Popen(
                    ["cmd.exe", "/c", "start", "", start_target, url],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                )
                return normalized, url
            except Exception as exc:  # pragma: no cover - best effort on host desktop
                last_error = exc
            continue
        if re.match(r"^[a-z][a-z0-9+.-]*:$", lower_command):
            try:
                subprocess.Popen(
                    ["cmd.exe", "/c", "start", "", f"{command}{url}"],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                )
                return normalized, url
            except Exception as exc:  # pragma: no cover - best effort on host desktop
                last_error = exc
            continue
        try:
            resolved = shutil.which(command) or command
            subprocess.Popen([resolved, url], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            return normalized, url
        except Exception as exc:  # pragma: no cover - best effort on host desktop
            last_error = exc

    friendly = START_MENU_FRIENDLY_NAMES.get(normalized, app_name.strip())
    raise RuntimeError(f"Could not open website '{url}' in '{friendly}': {last_error}")


def _resolve_folder_path(target: str | None, payload: dict | None = None) -> Path:
    payload = payload or {}
    raw_target = (target or payload.get("path") or payload.get("source_path") or "").strip().strip('"')
    if raw_target:
        candidate = Path(raw_target).expanduser()
        return candidate

    folder_key = str(payload.get("folder_key", "")).strip().lower()
    if folder_key in SPECIAL_FOLDERS:
        return SPECIAL_FOLDERS[folder_key]

    raise ValueError("A folder path is required.")


def _create_folder(base_path: Path, folder_name: str) -> Path:
    if not folder_name.strip():
        raise ValueError("A folder name is required.")
    created = base_path / folder_name.strip()
    created.mkdir(parents=True, exist_ok=True)
    return created


def _convert_pdfs_in_folder_to_ppts(source_path: Path, output_folder_name: str) -> dict[str, str | int]:
    if not source_path.exists() or not source_path.is_dir():
        raise ValueError(f"Source folder does not exist: {source_path}")

    output_dir = _create_folder(source_path, output_folder_name or "Converted_PPTs")
    pdf_files = sorted(source_path.glob("*.pdf"))
    if not pdf_files:
        raise ValueError(f"No PDF files found in {source_path}")

    converted_count = 0

    for pdf_file in pdf_files:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        with fitz.open(pdf_file) as document:
            for page_index in range(document.page_count):
                page = document.load_page(page_index)
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_image:
                    temp_path = Path(temp_image.name)
                try:
                    pixmap.save(temp_path.as_posix())
                    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                    slide.shapes.add_picture(
                        str(temp_path),
                        0,
                        0,
                        width=presentation.slide_width,
                        height=presentation.slide_height,
                    )
                finally:
                    if temp_path.exists():
                        temp_path.unlink(missing_ok=True)

        output_file = output_dir / f"{pdf_file.stem}.pptx"
        presentation.save(str(output_file))
        converted_count += 1

    return {
        "output_dir": str(output_dir),
        "converted_count": converted_count,
    }


def _open_external_url(url: str) -> str:
    normalized = url.strip()
    if not normalized:
        raise ValueError("A URL or search target is required.")
    if not normalized.startswith(("http://", "https://")):
        normalized = f"https://{normalized}"

    if hasattr(os, "startfile"):
        os.startfile(normalized)  # type: ignore[attr-defined]
    else:
        import webbrowser

        webbrowser.open(normalized, new=2)
    return normalized


def _extract_amount(payload: dict, default: int = 5) -> int:
    raw_value = payload.get("amount", payload.get("steps", payload.get("value", default)))
    try:
        parsed = int(float(str(raw_value)))
    except (TypeError, ValueError):
        parsed = default
    return max(1, min(parsed, 100))


def _set_windows_brightness(value: int) -> int:
    safe_value = max(0, min(int(value), 100))
    command = (
        "$monitor = Get-CimInstance -Namespace root/WMI -ClassName WmiMonitorBrightnessMethods "
        "-ErrorAction Stop | Select-Object -First 1; "
        f"Invoke-CimMethod -InputObject $monitor -MethodName WmiSetBrightness "
        f"-Arguments @{{Timeout=1;Brightness={safe_value}}} | Out-Null"
    )
    completed = subprocess.run(
        ["powershell.exe", "-NoProfile", "-Command", command],
        capture_output=True,
        text=True,
        timeout=15,
    )
    if completed.returncode != 0:
        raise RuntimeError(completed.stderr.strip() or "Windows brightness control is unavailable on this display.")
    return safe_value


def _get_windows_brightness() -> int:
    command = (
        "(Get-CimInstance -Namespace root/WMI -ClassName WmiMonitorBrightness "
        "-ErrorAction Stop | Select-Object -First 1).CurrentBrightness"
    )
    completed = subprocess.run(
        ["powershell.exe", "-NoProfile", "-Command", command],
        capture_output=True,
        text=True,
        timeout=15,
    )
    if completed.returncode != 0:
        raise RuntimeError(completed.stderr.strip() or "Windows brightness status is unavailable on this display.")
    match = re.search(r"\d+", completed.stdout)
    if not match:
        raise RuntimeError("Windows did not return a current brightness value.")
    return max(0, min(int(match.group(0)), 100))


async def execute_desktop_command(action: str, target: str = None, payload: dict | None = None):
    """Executes desktop and browser automation commands with best-effort safety."""
    payload = payload or {}

    try:
        if action == "open_notepad":
            pyautogui.press("win")
            await asyncio.sleep(0.5)
            pyautogui.write("notepad")
            await asyncio.sleep(0.5)
            pyautogui.press("enter")
            return {"success": True, "message": "Opened Notepad."}

        if action == "open_app":
            app_name = target or payload.get("app_name", "")
            opened_app = _launch_windows_app(app_name)
            return {
                "success": True,
                "message": f"Opened {opened_app}.",
                "note": "Window focus and background behavior still depend on the operating system.",
            }

        if action == "open_app_url":
            app_name = target or payload.get("app_name", "")
            target_url = payload.get("url", "")
            opened_app, normalized_url = _launch_windows_app_with_url(app_name, target_url)
            return {
                "success": True,
                "message": f"Opened {normalized_url} in {opened_app}.",
                "note": "Window focus and background behavior still depend on the operating system.",
            }

        if action == "whatsapp_send_message":
            contact = str(payload.get("contact", target or "")).strip()
            message_text = str(payload.get("message", payload.get("text", ""))).strip()
            if message_text:
                return await _whatsapp_open_contact_and_send(contact, message_text)
            return await _whatsapp_open_contact(contact)

        if action == "open_path":
            path = _resolve_folder_path(target, payload)
            if not path.exists():
                return {"success": False, "message": f"Path does not exist: {path}"}
            os.startfile(str(path))  # type: ignore[attr-defined]
            return {
                "success": True,
                "message": f"Opened {path} in File Explorer.",
                "note": "The operating system decides whether Explorer steals focus.",
            }

        if action == "create_folder":
            base_path = _resolve_folder_path(None, payload)
            folder_name = payload.get("folder_name", "")
            created_path = _create_folder(base_path, folder_name)
            return {
                "success": True,
                "message": f"Created folder {created_path.name}.",
                "path": str(created_path),
            }

        if action == "run_command":
            command = target or payload.get("command", "")
            shell_name = str(payload.get("shell", "powershell")).lower()
            if not command:
                return {"success": False, "message": "No command was provided."}

            if shell_name == "cmd":
                completed = subprocess.run(
                    ["cmd.exe", "/c", command],
                    capture_output=True,
                    text=True,
                    timeout=120,
                )
            else:
                completed = subprocess.run(
                    ["powershell.exe", "-NoProfile", "-Command", command],
                    capture_output=True,
                    text=True,
                    timeout=120,
                )

            return {
                "success": completed.returncode == 0,
                "message": "Ran the requested command." if completed.returncode == 0 else "The command failed.",
                "stdout": completed.stdout[-4000:],
                "stderr": completed.stderr[-4000:],
                "returncode": completed.returncode,
            }

        if action == "convert_pdfs_to_ppts":
            source_path = _resolve_folder_path(None, payload)
            output_folder_name = str(payload.get("output_folder_name", "Converted_PPTs")).strip() or "Converted_PPTs"
            result = _convert_pdfs_in_folder_to_ppts(source_path, output_folder_name)
            return {
                "success": True,
                "message": f"Converted {result['converted_count']} PDF file(s) into PPTX files.",
                "output_dir": result["output_dir"],
            }

        if action == "unsupported_browser_workflow":
            return {
                "success": False,
                "message": payload.get(
                    "message",
                    "This website workflow needs a DOM-aware web agent instead of active-field typing.",
                ),
                "note": payload.get("note"),
            }

        if action in {"type", "type_text"}:
            text = target or payload.get("text")
            if not text:
                return {"success": False, "message": "No text was provided to type."}
            pyautogui.write(text, interval=0.03)
            return {"success": True, "message": f"Typed text into the active window."}

        if action == "type_sequence":
            values = payload.get("values", [])
            if not values:
                return {"success": False, "message": "No values were provided for the fill sequence."}
            tab_presses_before = int(payload.get("tab_presses_before", 0))
            clear_each = bool(payload.get("clear_each", True))
            type_interval = float(payload.get("type_interval", 0.05))
            step_delay = float(payload.get("step_delay", 0.12))
            submit_delay = float(payload.get("submit_delay", 0.2))

            for _ in range(max(0, tab_presses_before)):
                pyautogui.press("tab")
                await asyncio.sleep(step_delay)

            for index, value in enumerate(values):
                if clear_each:
                    pyautogui.hotkey("ctrl", "a")
                    await asyncio.sleep(step_delay)
                    pyautogui.press("backspace")
                    await asyncio.sleep(step_delay)
                if value:
                    pyautogui.write(str(value), interval=max(0.02, type_interval))
                if index < len(values) - 1:
                    pyautogui.press("tab")
                    await asyncio.sleep(step_delay)
            if payload.get("submit"):
                await asyncio.sleep(submit_delay)
                pyautogui.press("enter")
            return {"success": True, "message": "Typed the form sequence into the active browser fields."}

        if action == "wait":
            delay = float(payload.get("seconds", 1.0))
            await asyncio.sleep(max(0.1, min(delay, 10)))
            return {"success": True, "message": f"Waited for {delay:.1f} seconds."}

        if action == "open_url":
            opened_url = _open_external_url(target or payload.get("url", ""))
            return {
                "success": True,
                "message": f"Opened {opened_url} in the default browser.",
                "note": "The browser decides whether the new page steals focus or opens in the background.",
            }

        if action == "open_youtube_song":
            query = target or payload.get("query", "")
            if not query:
                return {"success": False, "message": "Add a song or search query first."}
            opened_url = _open_external_url(
                f"https://www.youtube.com/results?search_query={quote_plus(query)}"
            )
            return {
                "success": True,
                "message": "Opened the YouTube search in the default browser.",
                "url": opened_url,
                "note": "Playing media in the background still depends on the browser and site autoplay rules.",
            }

        if action == "volume_up":
            presses = _extract_amount(payload, 5)
            pyautogui.press("volumeup", presses=presses, interval=0.03)
            return {"success": True, "message": f"Increased system volume by {presses} step(s)."}

        if action == "volume_down":
            presses = _extract_amount(payload, 5)
            pyautogui.press("volumedown", presses=presses, interval=0.03)
            return {"success": True, "message": f"Decreased system volume by {presses} step(s)."}

        if action in {"volume_mute", "volume_unmute"}:
            pyautogui.press("volumemute")
            return {"success": True, "message": "Toggled system mute."}

        if action == "set_brightness":
            value = _extract_amount(payload, 60)
            applied = _set_windows_brightness(value)
            return {"success": True, "message": f"Set screen brightness to {applied}%."}

        if action in {"brightness_up", "brightness_down"}:
            amount = _extract_amount(payload, 10)
            current = _get_windows_brightness()
            target_value = current + amount if action == "brightness_up" else current - amount
            applied = _set_windows_brightness(target_value)
            verb = "Increased" if action == "brightness_up" else "Decreased"
            return {"success": True, "message": f"{verb} screen brightness to {applied}%."}

        if action == "new_tab":
            pyautogui.hotkey("ctrl", "t")
            return {"success": True, "message": "Opened a new tab in the active browser window."}

        if action == "close_tab":
            pyautogui.hotkey("ctrl", "w")
            return {"success": True, "message": "Closed the current browser tab."}

        if action == "close_window":
            pyautogui.hotkey("alt", "f4")
            return {"success": True, "message": "Closed the current desktop window."}

        if action == "switch_window":
            pyautogui.hotkey("alt", "tab")
            return {"success": True, "message": "Switched to the next desktop window."}

        if action == "minimize_window":
            pyautogui.hotkey("win", "down")
            return {"success": True, "message": "Minimized the current window."}

        if action == "maximize_window":
            pyautogui.hotkey("win", "up")
            return {"success": True, "message": "Maximized the current window."}

        if action == "edit_field":
            text = target or payload.get("text", "")
            pyautogui.hotkey("ctrl", "a")
            await asyncio.sleep(0.08)
            if text:
                pyautogui.write(text, interval=0.03)
            return {"success": True, "message": "Replaced the active field content."}

        if action in {"clear_field", "remove_draft"}:
            pyautogui.hotkey("ctrl", "a")
            await asyncio.sleep(0.08)
            pyautogui.press("backspace")
            return {"success": True, "message": "Cleared the active field or draft content."}

        return {
            "success": False,
            "message": f"Automation action '{action}' is not configured yet.",
        }
    except Exception as e:
        return {"success": False, "message": f"Automation failed: {str(e)}"}
